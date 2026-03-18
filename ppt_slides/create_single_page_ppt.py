from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9
prs.slide_height = Inches(7.5)

# Define colors - Professional dark theme
DARK_BG = RGBColor(15, 23, 42)      # #0f172a
LIGHT_BG = RGBColor(30, 41, 59)     # #1e293b
BLUE = RGBColor(59, 130, 246)       # #3b82f6
GREEN = RGBColor(34, 197, 94)       # #22c55e
YELLOW = RGBColor(251, 191, 36)     # #fbbf24
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(148, 163, 184)      # #94a3b8
DARK_GRAY = RGBColor(100, 116, 139) # #64748b

# Create blank slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = DARK_BG

# Title at top
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "PairProxy - 企业级 LLM API 代理网关"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Subtitle
subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.333), Inches(0.4))
tf = subtitle_box.text_frame
p = tf.paragraphs[0]
p.text = "统一管控 · 精确追踪 · 零侵入接入 · AI 驱动开发"
p.font.size = Pt(14)
p.font.color.rgb = BLUE
p.alignment = PP_ALIGN.CENTER

# Layout: 3 columns
# Column 1: Core Stats & Architecture
col1_x = 0.5
col_width = 4
row_height = 0.45

# Column 1 Title
col1_title = slide.shapes.add_textbox(Inches(col1_x), Inches(1.6), Inches(col_width), Inches(0.4))
tf = col1_title.text_frame
p = tf.paragraphs[0]
p.text = "📊 核心数据"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = YELLOW

# Stats in column 1
stats = [
    ("64,000+", "Go 代码行数"),
    ("1,100+", "测试用例 全部通过"),
    ("~70%", "测试覆盖率"),
    ("v2.5.0", "Production Ready")
]

y_pos = 2.0
for num, desc in stats:
    # Number
    num_box = slide.shapes.add_textbox(Inches(col1_x), Inches(y_pos), Inches(1.5), Inches(0.35))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN

    # Description
    desc_box = slide.shapes.add_textbox(Inches(col1_x + 1.6), Inches(y_pos), Inches(2.4), Inches(0.35))
    tf = desc_box.text_frame
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY

    y_pos += 0.45

# Column 2: Key Features
col2_x = 5.0

col2_title = slide.shapes.add_textbox(Inches(col2_x), Inches(1.6), Inches(col_width), Inches(0.4))
tf = col2_title.text_frame
p = tf.paragraphs[0]
p.text = "⚡ 核心能力"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = YELLOW

features = [
    "JWT 认证 + LDAP/AD 集成",
    "日/月 Token 配额 + RPM 限流",
    "两级负载均衡 + 健康检查",
    "实时费用统计 + Dashboard",
    "对话内容追踪 + 审计日志",
    "集群模式 + 自动故障恢复"
]

y_pos = 2.0
for feature in features:
    feat_box = slide.shapes.add_textbox(Inches(col2_x), Inches(y_pos), Inches(col_width), Inches(0.35))
    tf = feat_box.text_frame
    p = tf.paragraphs[0]
    p.text = "• " + feature
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    y_pos += 0.42

# Column 3: Business Value
col3_x = 9.5

col3_title = slide.shapes.add_textbox(Inches(col3_x), Inches(1.6), Inches(col_width), Inches(0.4))
tf = col3_title.text_frame
p = tf.paragraphs[0]
p.text = "💼 业务价值"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = YELLOW

values = [
    ("成本透明", "精确到 0.000001 USD"),
    ("预算可控", "超额自动拦截 (429)"),
    ("责任可追", "完整审计日志"),
    ("安全合规", "API Key 集中管控")
]

y_pos = 2.0
for title, desc in values:
    # Title
    val_title = slide.shapes.add_textbox(Inches(col3_x), Inches(y_pos), Inches(col_width), Inches(0.3))
    tf = val_title.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = GREEN

    # Description
    val_desc = slide.shapes.add_textbox(Inches(col3_x), Inches(y_pos + 0.28), Inches(col_width), Inches(0.3))
    tf = val_desc.text_frame
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY

    y_pos += 0.65

# Architecture diagram at bottom
arch_y = 4.8

# Architecture section title
arch_title = slide.shapes.add_textbox(Inches(0.5), Inches(arch_y), Inches(12.333), Inches(0.35))
tf = arch_title.text_frame
p = tf.paragraphs[0]
p.text = "🔄 系统架构"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = YELLOW
p.alignment = PP_ALIGN.CENTER

# Architecture flow
arch_box = slide.shapes.add_textbox(Inches(0.5), Inches(arch_y + 0.4), Inches(12.333), Inches(0.5))
tf = arch_box.text_frame
p = tf.paragraphs[0]
p.text = "Claude Code → cproxy(本地:8080) → sproxy(服务端:9000) → Anthropic / OpenAI / Ollama"
p.font.size = Pt(13)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Bottom stats bar
bar_y = 6.0
stats_bar = [
    ("99.5%", "可用性"),
    ("4.8/5.0", "成熟度评分"),
    ("50+", "团队规模支持"),
    ("Apache 2.0", "开源协议")
]

stat_width = 3.0
start_x = (13.333 - len(stats_bar) * stat_width) / 2

for i, (num, label) in enumerate(stats_bar):
    x = start_x + i * stat_width

    # Number
    num_box = slide.shapes.add_textbox(Inches(x), Inches(bar_y), Inches(stat_width), Inches(0.5))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p.alignment = PP_ALIGN.CENTER

    # Label
    lbl_box = slide.shapes.add_textbox(Inches(x), Inches(bar_y + 0.5), Inches(stat_width), Inches(0.35))
    tf = lbl_box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

# Footer
footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.333), Inches(0.4))
tf = footer_box.text_frame
p = tf.paragraphs[0]
p.text = "github.com/l17728/pairproxy | AI Coding 开发企业级软件的成功实践"
p.font.size = Pt(10)
p.font.color.rgb = DARK_GRAY
p.alignment = PP_ALIGN.CENTER

# Save
prs.save('D:\pairproxy\PairProxy_OnePage_Summary.pptx')
print("PPT created: D:\pairproxy\PairProxy_OnePage_Summary.pptx")
print("Total slides: 1 (浓缩版)")
